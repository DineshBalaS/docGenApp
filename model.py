from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptInches
import os
import re
from PIL import Image

TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), 'templates_src')
DOCX_TEMPLATE = os.path.join(TEMPLATE_DIR, 'template.docx')
PPTX_TEMPLATE = os.path.join(TEMPLATE_DIR, 'newPpt.pptx')

IMAGE_FIELDS = [
    "property - sample 1 - screenshot",
    "property - sample 1 - image",
    "property - sample 2 - screenshot",
    "property - sample 2 - image",
    "Property image 1",
    "Property image 2",
    "Property 3D images 1",
    "Property 3D images 2",
    "Interior image 1",
    "Interior image 2",
    "Interior image 3",
    "Interior image 4",
    "OS Map",
    "Sim app Floor plan 1",
    "Sim app elevation 1",
    "Sim app floor plan 2",
    "Sim app elevation 2",
    "rejected application- image",
    "Conceptual Design Chatgpt",
    "Conceptual Floor Plan",
    "scope of work"
]

if not os.path.exists(DOCX_TEMPLATE):
    raise FileNotFoundError(f"DOCX template not found at: {DOCX_TEMPLATE}")

# NEW, SIMPLIFIED FUNCTION: Replaces old clean_placeholder and find_matching_key
def find_matching_key(raw_placeholder, data_keys):
    """
    Finds a matching key in the data dictionary by comparing lowercase, stripped versions.
    """
    raw_placeholder_clean = raw_placeholder.strip().lower()
    for key in data_keys:
        if key.lower() == raw_placeholder_clean:
            return key
    return None

def replace_in_paragraphs(paragraphs, data):
    # This function is for the DOCX generator and seems okay.
    # For simplicity, we are focusing on the PPTX generator.
    for para in paragraphs:
        # NOTE: This simple text replace will not handle fragmented runs in DOCX
        matches = re.findall(r"\{\{(.*?)\}\}", para.text)
        for raw_ph in matches:
            actual_key = find_matching_key(raw_ph, data)
            if actual_key and actual_key not in IMAGE_FIELDS:
                value = data.get(actual_key, "")
                if isinstance(value, list): value = ", ".join(value)
                para.text = para.text.replace(f"{{{{{raw_ph}}}}}", str(value))

def replace_in_tables(tables, data):
    for table in tables:
        for row in table.rows:
            for cell in row.cells:
                replace_in_paragraphs(cell.paragraphs, data)

def generate_docx(data: dict, output_path: str):
    doc = Document(DOCX_TEMPLATE)
    replace_in_paragraphs(doc.paragraphs, data)
    replace_in_tables(doc.tables, data)
    doc.save(output_path)

def generate_pptx(data: dict, output_path: str):
    ppt = Presentation(PPTX_TEMPLATE)
    for slide_idx, slide in enumerate(ppt.slides):
        shapes_to_delete = []

        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue

            full_shape_text = "".join(run.text for p in shape.text_frame.paragraphs for run in p.runs)
            all_matches = re.findall(r"\{\{(.*?)\}\}", full_shape_text)
            
            if all_matches:
                print(f"\n--- Placeholders Found on Slide {slide_idx + 1} ---")
                for ph in all_matches:
                    key_found = find_matching_key(ph, data.keys())
                    print(f"  - Reading '{ph}', Matched Key: {key_found}")

            if not all_matches:
                continue

            is_image_shape = False
            image_key = None
            image_path = None

            for raw_ph in all_matches:
                actual_key = find_matching_key(raw_ph, data.keys())
                if actual_key and actual_key in IMAGE_FIELDS:
                    is_image_shape = True
                    image_key = actual_key
                    value = data.get(image_key, "")
                    image_path = os.path.abspath(value) if value else ""
                    break
            
            if is_image_shape:
                if image_path and os.path.isfile(image_path):
                    try:
                        slide.shapes.add_picture(image_path, shape.left, shape.top, width=shape.width, height=shape.height)
                        shapes_to_delete.append(shape)
                    except Exception as e:
                        print(f"ERROR: Could not add image {image_path}. Details: {e}")
                        shape.text_frame.text = f"[Image Error: {image_key}]"
                else:
                    shape.text_frame.text = f"[Image Missing: {image_key}]"
                continue

            if "{{list of documents required}}" in full_shape_text:
                tf = shape.text_frame
                tf.clear()
                doc_list = data.get('list of documents required', [])
                if doc_list:
                    first_p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                    first_p.text = doc_list[0]
                    first_p.level = 0
                    for doc in doc_list[1:]:
                        p = tf.add_paragraph()
                        p.text = doc
                        p.level = 0
                continue
            
            for para in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in para.runs)
                para_matches = re.findall(r"\{\{(.*?)\}\}", para_text)

                if not para_matches:
                    continue
                
                final_text = para_text

                for raw_ph in para_matches:
                    actual_key = find_matching_key(raw_ph, data.keys())
                    if actual_key:
                        value = data.get(actual_key, f"[{actual_key} ?]")
                        if isinstance(value, list): value = ", ".join(value)
                        final_text = final_text.replace(f"{{{{{raw_ph}}}}}", str(value))
                
                para.clear()
                run = para.add_run()
                run.text = final_text
        
        for shape in shapes_to_delete:
            sp_element = shape.element
            sp_element.getparent().remove(sp_element)

    ppt.save(output_path)